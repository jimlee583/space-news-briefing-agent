from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from space_news_briefing_agent.deck import build_deck
from space_news_briefing_agent.models import Briefing, CompanyBrief, NewsItem


def _briefing() -> Briefing:
    item = NewsItem(
        title="Rocket Lab books new SDA contract",
        source="SpaceNews",
        date="2026-05-01",
        url="https://example.com/rocketlab-sda",
        summary="Rocket Lab announced a new SDA tranche award worth $X.",
        why_it_matters="Reinforces small-sat primes' position in proliferated LEO architectures.",
        confidence="high",
    )
    return Briefing(
        deck_title="Daily Space & Defense-Space News Briefing",
        briefing_date="2026-05-01",
        cross_company_summary="Mostly quiet news day; one notable Rocket Lab award.",
        company_briefs=[
            CompanyBrief(
                topic_name="Rocket Lab",
                executive_summary="Rocket Lab landed a new SDA award.",
                top_items=[item],
                implications=["Strengthens proliferated LEO posture."],
                watch_items=["Watch for Neutron schedule updates."],
            ),
            CompanyBrief(
                topic_name="K2 Space",
                executive_summary="No qualifying news in the lookback window.",
                top_items=[],
                implications=[],
                watch_items=[],
            ),
        ],
        industry_themes=["SDA tranche cadence is accelerating."],
        defense_space_implications=["More resilient missile-warning architectures."],
        watch_items=["Upcoming SSC awards."],
        source_list=["https://example.com/rocketlab-sda"],
    )


def test_build_deck_writes_file_with_expected_slides(tmp_path: Path) -> None:
    deck_path = build_deck(_briefing(), tmp_path)

    assert deck_path.exists()
    assert deck_path.name == "space_defense_news_briefing_2026-05-01.pptx"

    prs = Presentation(deck_path)
    titles = []
    for slide in prs.slides:
        title = slide.shapes.title.text if slide.shapes.title is not None else ""
        titles.append(title)

    assert "Daily Space & Defense-Space News Briefing" in titles[0]
    assert "Executive Summary" in titles
    assert "Industry Themes" in titles
    assert "Defense-Space Implications" in titles
    assert any("Rocket Lab — Summary" in t for t in titles)
    assert any("K2 Space — Summary" in t for t in titles)
    assert any("Rocket Lab — Rocket Lab books" in t for t in titles)
    assert "Watch Items" in titles
    assert "Sources" in titles
