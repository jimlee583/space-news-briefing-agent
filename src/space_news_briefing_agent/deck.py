"""PowerPoint deck generator.

The deck layout is intentionally simple and theme-agnostic so it works with
any default `python-pptx` template. Slides are added in this order:

1. Title
2. Cross-company executive summary
3. News highlights (industry themes)
4. Upcoming launches & schedule watch (when launch events are present)
5. Defense-space implications
6. One summary slide per company / topic
7. One slide per top news item, grouped by company
8. Watch items
9. Sources
"""

from __future__ import annotations

import logging
from collections.abc import Iterable
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from .core.models import IntelligenceEvent, LaunchEvent
from .models import Briefing, CompanyBrief, NewsItem

logger = logging.getLogger(__name__)

# Layout indices for the default `python-pptx` template:
#   0: Title Slide, 1: Title and Content, 5: Title Only, 6: Blank
_LAYOUT_TITLE = 0
_LAYOUT_TITLE_AND_CONTENT = 1
_LAYOUT_TITLE_ONLY = 5

_BODY_FONT_SIZE_PT = 16
_BULLET_FONT_SIZE_PT = 14
_FOOTER_FONT_SIZE_PT = 10


# --------------------------------------------------------------------------- #
# Public API
# --------------------------------------------------------------------------- #


def build_deck(
    briefing: Briefing,
    output_dir: Path,
    *,
    events: Iterable[IntelligenceEvent] | None = None,
) -> Path:
    """Render `briefing` to `output_dir/space_defense_news_briefing_YYYY-MM-DD.pptx`.

    ``events`` is optional and used only to render the "Upcoming Launches"
    slide. Passing ``None`` preserves the legacy news-only deck.

    Returns the absolute path to the generated file.
    """
    output_dir = output_dir.expanduser().resolve()
    output_dir.mkdir(parents=True, exist_ok=True)

    launch_events: list[LaunchEvent] = (
        [e for e in events if isinstance(e, LaunchEvent)] if events else []
    )

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _add_title_slide(prs, briefing)
    _add_summary_slide(
        prs,
        title="Executive Summary",
        body_paragraphs=[briefing.cross_company_summary]
        if briefing.cross_company_summary
        else ["No major updates found."],
    )
    _add_bulleted_slide(prs, "News Highlights", briefing.industry_themes)

    if launch_events:
        _add_launch_slide(prs, launch_events)

    _add_bulleted_slide(prs, "Defense-Space Implications", briefing.defense_space_implications)

    for brief in briefing.company_briefs:
        _add_company_summary_slide(prs, brief)

    for brief in briefing.company_briefs:
        for item in brief.top_items:
            _add_news_item_slide(prs, brief.topic_name, item)

    _add_bulleted_slide(prs, "Watch Items", briefing.watch_items)
    _add_sources_slide(prs, [str(u) for u in briefing.source_list])

    filename = f"space_defense_news_briefing_{briefing.briefing_date}.pptx"
    output_path = output_dir / filename
    prs.save(output_path)
    logger.info("Wrote deck: %s", output_path)
    return output_path


# --------------------------------------------------------------------------- #
# Slide builders
# --------------------------------------------------------------------------- #


def _add_title_slide(prs: Presentation, briefing: Briefing) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[_LAYOUT_TITLE])
    title = slide.shapes.title
    if title is not None:
        title.text = briefing.deck_title or "Daily Space & Defense-Space News Briefing"
    if len(slide.placeholders) > 1:
        subtitle = slide.placeholders[1]
        subtitle.text = f"Briefing date: {briefing.briefing_date}"


def _add_summary_slide(prs: Presentation, *, title: str, body_paragraphs: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[_LAYOUT_TITLE_AND_CONTENT])
    _set_title(slide, title)
    body = _content_placeholder(slide)
    if body is None:
        return
    tf = body.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, paragraph in enumerate(body_paragraphs):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = paragraph
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(_BODY_FONT_SIZE_PT)


def _add_bulleted_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[_LAYOUT_TITLE_AND_CONTENT])
    _set_title(slide, title)
    body = _content_placeholder(slide)
    if body is None:
        return
    tf = body.text_frame
    tf.word_wrap = True
    tf.clear()

    items = [b.strip() for b in bullets if b and b.strip()]
    if not items:
        items = ["No items reported."]

    for idx, bullet in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(_BULLET_FONT_SIZE_PT)


def _add_company_summary_slide(prs: Presentation, brief: CompanyBrief) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[_LAYOUT_TITLE_AND_CONTENT])
    _set_title(slide, f"{brief.topic_name} — Summary")
    body = _content_placeholder(slide)
    if body is None:
        return
    tf = body.text_frame
    tf.word_wrap = True
    tf.clear()

    p = tf.paragraphs[0]
    p.text = brief.executive_summary or "No summary available."
    for run in p.runs:
        run.font.size = Pt(_BODY_FONT_SIZE_PT)

    if brief.implications:
        _add_section_heading(tf, "Implications")
        for line in brief.implications:
            _add_bullet(tf, line, level=1)

    if brief.watch_items:
        _add_section_heading(tf, "Watch")
        for line in brief.watch_items:
            _add_bullet(tf, line, level=1)


def _add_news_item_slide(prs: Presentation, topic_name: str, item: NewsItem) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[_LAYOUT_TITLE_AND_CONTENT])
    _set_title(slide, f"{topic_name} — {_truncate(item.title, 90)}")
    body = _content_placeholder(slide)
    if body is None:
        return
    tf = body.text_frame
    tf.word_wrap = True
    tf.clear()

    meta_bits = [item.source]
    if item.date:
        meta_bits.append(item.date)
    meta_bits.append(f"confidence: {item.confidence}")
    meta = " · ".join(b for b in meta_bits if b)

    p = tf.paragraphs[0]
    p.text = meta
    for run in p.runs:
        run.font.size = Pt(_FOOTER_FONT_SIZE_PT)
        run.font.italic = True

    _add_section_heading(tf, "Summary")
    _add_bullet(tf, item.summary, level=1)

    _add_section_heading(tf, "Why it matters")
    _add_bullet(tf, item.why_it_matters, level=1)

    _add_section_heading(tf, "Source")
    _add_bullet(tf, str(item.url), level=1)


def _add_launch_slide(prs: Presentation, launches: list[LaunchEvent]) -> None:
    """Render a "Upcoming Launches & Schedule Watch" slide.

    Sorted by launch T-0 (earliest first). Limits to the first 12 entries to
    keep the slide readable; full data lives in ``output/events.jsonl``.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[_LAYOUT_TITLE_AND_CONTENT])
    _set_title(slide, "Upcoming Launches & Schedule Watch")
    body = _content_placeholder(slide)
    if body is None:
        return
    tf = body.text_frame
    tf.word_wrap = True
    tf.clear()

    if not launches:
        p = tf.paragraphs[0]
        p.text = "No upcoming launches relevant to tracked entities in the lookahead window."
        for run in p.runs:
            run.font.size = Pt(_BULLET_FONT_SIZE_PT)
        return

    ordered = sorted(launches, key=_launch_sort_key)[:12]
    for idx, launch in enumerate(ordered):
        line = _format_launch_line(launch)
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(_BULLET_FONT_SIZE_PT)


def _launch_sort_key(launch: LaunchEvent) -> tuple[int, str]:
    if launch.event_date is None:
        return (1, launch.title.lower())
    return (0, launch.event_date.isoformat())


def _format_launch_line(launch: LaunchEvent) -> str:
    when = launch.event_date.strftime("%Y-%m-%d %H:%MZ") if launch.event_date else "TBD"
    bits = [when]
    if launch.launch_provider:
        bits.append(launch.launch_provider)
    if launch.vehicle:
        bits.append(launch.vehicle)
    if launch.payload:
        bits.append(launch.payload)
    elif launch.title:
        bits.append(launch.title)
    if launch.launch_site:
        bits.append(launch.launch_site)
    if launch.mission_status:
        bits.append(f"[{launch.mission_status}]")
    return " · ".join(b for b in bits if b)


def _add_sources_slide(prs: Presentation, urls: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[_LAYOUT_TITLE_AND_CONTENT])
    _set_title(slide, "Sources")
    body = _content_placeholder(slide)
    if body is None:
        return
    tf = body.text_frame
    tf.word_wrap = True
    tf.clear()

    if not urls:
        p = tf.paragraphs[0]
        p.text = "No sources cited."
        for run in p.runs:
            run.font.size = Pt(_BULLET_FONT_SIZE_PT)
        return

    for idx, url in enumerate(urls):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = url
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(_FOOTER_FONT_SIZE_PT)


# --------------------------------------------------------------------------- #
# Low-level helpers
# --------------------------------------------------------------------------- #


def _set_title(slide: object, text: str) -> None:
    title = getattr(slide.shapes, "title", None)  # type: ignore[attr-defined]
    if title is not None:
        title.text = text


def _content_placeholder(slide: object):  # type: ignore[no-untyped-def]
    """Return the first non-title placeholder, or None.

    `python-pptx` exposes placeholders by index; idx 0 is the title and idx 1
    is the body for the standard "Title and Content" layout.
    """
    placeholders = slide.placeholders  # type: ignore[attr-defined]
    for placeholder in placeholders:
        if placeholder.placeholder_format.idx != 0:
            return placeholder
    return None


def _add_section_heading(tf, text: str) -> None:  # type: ignore[no-untyped-def]
    p = tf.add_paragraph()
    p.text = text
    p.level = 0
    for run in p.runs:
        run.font.size = Pt(_BULLET_FONT_SIZE_PT)
        run.font.bold = True


def _add_bullet(tf, text: str, *, level: int) -> None:  # type: ignore[no-untyped-def]
    p = tf.add_paragraph()
    p.text = text
    p.level = level
    for run in p.runs:
        run.font.size = Pt(_BULLET_FONT_SIZE_PT)


def _truncate(text: str, limit: int) -> str:
    if len(text) <= limit:
        return text
    return text[: limit - 1].rstrip() + "…"


# Silence unused-import warnings for shapes constants kept for future styling hooks.
_ = MSO_SHAPE
